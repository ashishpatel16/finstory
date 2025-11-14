"""
PowerPoint Presentation Generator
Creates beautiful PPTX presentations from JSON input
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from typing import List, Dict, Any, Optional
from dataclasses import dataclass
import json


@dataclass
class PresentationTheme:
    """Centralized color & typography palette for the deck."""
    background_color: RGBColor
    heading_level_one_color: RGBColor
    heading_level_two_color: RGBColor
    text_color: RGBColor
    surface_color: RGBColor
    accent_color: RGBColor
    positive_color: RGBColor
    negative_color: RGBColor
    font_family_heading: str = "Montserrat"
    font_family_body: str = "Inter"


DEFAULT_THEME = PresentationTheme(
    background_color=RGBColor(245, 249, 255),        # Arctic mist
    heading_level_one_color=RGBColor(15, 76, 129),   # Deep sapphire
    heading_level_two_color=RGBColor(56, 189, 248),  # Sky cyan
    text_color=RGBColor(30, 41, 59),                 # Slate navy
    surface_color=RGBColor(255, 255, 255),           # Pure white cards
    accent_color=RGBColor(79, 70, 229),              # Indigo accent
    positive_color=RGBColor(34, 197, 94),            # Success green
    negative_color=RGBColor(239, 68, 68),            # Alert red
)


class SlideGenerator:
    """Generates PowerPoint presentations from JSON data"""
    
    def __init__(
        self,
        width: float = 10.0,
        height: float = 7.5,
        theme: Optional[PresentationTheme] = None,
    ):
        """
        Initialize the presentation generator
        
        Args:
            width: Slide width in inches (default: 10.0 for widescreen)
            height: Slide height in inches (default: 7.5 for widescreen)
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(width)
        self.prs.slide_height = Inches(height)
        self.theme = theme or DEFAULT_THEME
    
    @staticmethod
    def _title_font_size(text: str) -> int:
        """Return a font size that keeps long titles inside the slide."""
        length = len(text)
        if length <= 20:
            return 54
        if length <= 35:
            return 46
        if length <= 55:
            return 40
        if length <= 80:
            return 34
        return 30
    
    def _add_title_slide(self, title: str, subtitle: Optional[str] = None):
        """Create a title slide with beautiful styling"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add background shape with theme background color
        left = top = Inches(0)
        width = self.prs.slide_width
        height = self.prs.slide_height
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme.background_color
        background.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.1), self.prs.slide_width - Inches(1.0), Inches(2.0)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(self._title_font_size(title))
        title_para.font.bold = True
        title_para.font.name = self.theme.font_family_heading
        title_para.font.color.rgb = self.theme.heading_level_one_color
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5), Inches(9), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.name = self.theme.font_family_body
            subtitle_para.font.color.rgb = self.theme.heading_level_two_color
            subtitle_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _add_content_slide(self, title: str, content: Any):
        """Create a content slide with various content types"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
            self.prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.theme.heading_level_one_color
        title_bar.line.fill.background()
        
        # Add title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.name = self.theme.font_family_heading
        title_para.font.color.rgb = self.theme.surface_color
        title_para.alignment = PP_ALIGN.LEFT
        
        # Add content based on type
        content_top = Inches(1.8)
        
        if isinstance(content, list):
            # Bullet points
            self._add_bullet_points(slide, content, content_top)
        elif isinstance(content, dict):
            # Structured content
            if 'type' in content:
                if content['type'] == 'text':
                    self._add_text_content(slide, content.get('text', ''), content_top)
                elif content['type'] == 'bullets':
                    self._add_bullet_points(slide, content.get('items', []), content_top)
                elif content['type'] == 'two_column':
                    self._add_two_column(slide, content.get('left', []), 
                                       content.get('right', []), content_top)
            else:
                # Key-value pairs
                self._add_key_value_pairs(slide, content, content_top)
        else:
            # Plain text
            self._add_text_content(slide, str(content), content_top)
        
        return slide
    
    def _add_bullet_points(self, slide, items: List[str], top: float):
        """Add bullet points to slide"""
        text_box = slide.shapes.add_textbox(
            Inches(0.8), top, Inches(8.4), Inches(5)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.level = 0
            p.font.size = Pt(20)
            p.font.name = self.theme.font_family_body
            p.font.color.rgb = self.theme.text_color
            p.space_after = Pt(12)
    
    def _add_text_content(self, slide, text: str, top: float):
        """Add plain text content to slide"""
        text_box = slide.shapes.add_textbox(
            Inches(0.8), top, Inches(8.4), Inches(5)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = text
        
        para = text_frame.paragraphs[0]
        para.font.size = Pt(20)
        para.font.name = self.theme.font_family_body
        para.font.color.rgb = self.theme.text_color
        para.alignment = PP_ALIGN.LEFT
    
    def _add_two_column(self, slide, left_items: List[str], right_items: List[str], top: float):
        """Add two-column layout"""
        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.8), top, Inches(4), Inches(5)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        for i, item in enumerate(left_items):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.name = self.theme.font_family_body
            p.font.color.rgb = self.theme.text_color
            p.space_after = Pt(10)
        
        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(5.2), top, Inches(4), Inches(5)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        for i, item in enumerate(right_items):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.name = self.theme.font_family_body
            p.font.color.rgb = self.theme.text_color
            p.space_after = Pt(10)
    
    def _add_key_value_pairs(self, slide, data: Dict[str, Any], top: float):
        """Add key-value pairs in a structured format"""
        text_box = slide.shapes.add_textbox(
            Inches(0.8), top, Inches(8.4), Inches(5)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, (key, value) in enumerate(data.items()):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = f"{key}: {value}"
            p.font.size = Pt(20)
            p.font.name = self.theme.font_family_body
            p.font.color.rgb = self.theme.text_color
            p.space_after = Pt(14)
            p.space_before = Pt(8)
    
    def generate_from_json(self, slides_data: List[Dict[str, Any]]) -> Presentation:
        """
        Generate presentation from JSON data
        
        Args:
            slides_data: List of slide dictionaries, each containing:
                - type: 'title' or 'content'
                - title: Slide title (required)
                - subtitle: Subtitle (for title slides)
                - content: Content (for content slides) - can be:
                    - String
                    - List of strings (bullet points)
                    - Dict with 'type' key ('text', 'bullets', 'two_column')
                    - Dict of key-value pairs
        
        Returns:
            Presentation object
        """
        for slide_data in slides_data:
            slide_type = slide_data.get('type', 'content')
            title = slide_data.get('title', 'Untitled')
            
            if slide_type == 'title':
                subtitle = slide_data.get('subtitle')
                self._add_title_slide(title, subtitle)
            else:
                content = slide_data.get('content', '')
                self._add_content_slide(title, content)
        
        return self.prs
    
    def _add_metric_card(
        self,
        slide,
        label: str,
        value: str,
        left: float,
        top: float,
        width: float,
        height: float,
        bg_color: Optional[RGBColor] = None,
        value_color: Optional[RGBColor] = None,
    ):
        """Add a beautiful metric card to slide"""
        # Background shape
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), 
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color or self.theme.surface_color
        card.line.color.rgb = self.theme.heading_level_two_color
        card.line.width = Pt(1)
        
        # Value text (large, bold)
        value_box = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.3), Inches(width - 0.4), Inches(0.8)
        )
        value_frame = value_box.text_frame
        value_frame.text = value
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(32)
        value_para.font.bold = True
        value_para.font.name = self.theme.font_family_heading
        value_para.font.color.rgb = value_color or self.theme.heading_level_one_color
        value_para.alignment = PP_ALIGN.CENTER
        
        # Label text (smaller)
        label_box = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 1.0), Inches(width - 0.4), Inches(0.5)
        )
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(14)
        label_para.font.name = self.theme.font_family_body
        label_para.font.color.rgb = self.theme.text_color
        label_para.alignment = PP_ALIGN.CENTER
    
    def _add_pnl_slide(self, title: str, pnl_data: Dict[str, Any]):
        """Create a Profit & Loss statement slide"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
            self.prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.theme.heading_level_one_color
        title_bar.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.name = self.theme.font_family_heading
        title_para.font.color.rgb = self.theme.surface_color
        title_para.alignment = PP_ALIGN.LEFT
        
        # P&L table
        top = Inches(1.8)
        row_height = Inches(0.5)
        left_col_width = Inches(5)
        right_col_width = Inches(4)
        
        # Header row
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), top, 
            left_col_width + right_col_width, row_height
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = self.theme.heading_level_two_color
        header_bg.line.fill.background()
        
        header_text = slide.shapes.add_textbox(
            Inches(0.6), top + Inches(0.1), left_col_width + right_col_width - Inches(0.2), row_height - Inches(0.2)
        )
        header_frame = header_text.text_frame
        header_frame.text = "Line Item"
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(16)
        header_para.font.bold = True
        header_para.font.name = self.theme.font_family_body
        header_para.font.color.rgb = self.theme.surface_color
        
        # P&L items
        items = [
            ("Revenue", pnl_data.get('revenue', 0), True),
            ("Cost of Goods Sold (COGS)", pnl_data.get('cogs', 0), False),
            ("Gross Profit", pnl_data.get('gross_profit', 0), True),
            ("Operating Expenses", pnl_data.get('opex', 0), False),
            ("EBITDA", pnl_data.get('ebitda', 0), True),
            ("Depreciation & Amortization", pnl_data.get('depreciation', 0), False),
            ("EBIT", pnl_data.get('ebit', 0), True),
            ("Interest Expense", pnl_data.get('interest', 0), False),
            ("Net Income", pnl_data.get('net_income', 0), True),
        ]
        
        def format_currency(value):
            if isinstance(value, (int, float)):
                if value >= 1_000_000:
                    return f"${value/1_000_000:.2f}M"
                elif value >= 1_000:
                    return f"${value/1_000:.2f}K"
                return f"${value:,.2f}"
            return str(value)
        
        current_top = top + row_height
        for i, (label, value, is_important) in enumerate(items):
            # Row background (alternating)
            if i % 2 == 0:
                row_bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0.5), current_top,
                    left_col_width + right_col_width, row_height
                )
                row_bg.fill.solid()
                row_bg.fill.fore_color.rgb = self.theme.surface_color
                row_bg.line.fill.background()
            
            # Label
            label_box = slide.shapes.add_textbox(
                Inches(0.6), current_top + Inches(0.1), left_col_width - Inches(0.2), row_height - Inches(0.2)
            )
            label_frame = label_box.text_frame
            label_frame.text = label
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(14 if not is_important else 16)
            label_para.font.bold = is_important
            label_para.font.name = self.theme.font_family_body
            label_para.font.color.rgb = self.theme.text_color
            
            # Value
            value_box = slide.shapes.add_textbox(
                Inches(5.5), current_top + Inches(0.1), right_col_width - Inches(0.2), row_height - Inches(0.2)
            )
            value_frame = value_box.text_frame
            value_frame.text = format_currency(value)
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(14 if not is_important else 18)
            value_para.font.bold = is_important
            value_para.font.name = self.theme.font_family_heading if is_important else self.theme.font_family_body
            value_para.font.color.rgb = (
                self.theme.accent_color if is_important else self.theme.text_color
            )
            value_para.alignment = PP_ALIGN.RIGHT
            
            current_top += row_height
    
    def _add_key_metrics_slide(self, title: str, metrics: Dict[str, Any]):
        """Create a key metrics overview slide with metric cards"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
            self.prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.theme.heading_level_one_color
        title_bar.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.name = self.theme.font_family_heading
        title_para.font.color.rgb = self.theme.surface_color
        title_para.alignment = PP_ALIGN.LEFT
        
        # Format currency helper
        def format_currency(value):
            if isinstance(value, (int, float)):
                if value >= 1_000_000:
                    return f"${value/1_000_000:.2f}M"
                elif value >= 1_000:
                    return f"${value/1_000:.2f}K"
                return f"${value:,.2f}"
            return str(value)
        
        def format_percent(value):
            if isinstance(value, (int, float)):
                return f"{value:.1f}%"
            return str(value)
        
        # Top row metrics
        top_y = 1.8
        card_width = 2.8
        card_height = 1.5
        spacing = 0.3
        
        # Revenue
        self._add_metric_card(
            slide,
            "Revenue",
            format_currency(metrics.get('revenue', 0)),
            0.5,
            top_y,
            card_width,
            card_height,
            value_color=self.theme.heading_level_one_color,
        )
        
        # Gross Margin
        self._add_metric_card(
            slide,
            "Gross Margin",
            format_percent(metrics.get('gross_margin', 0)),
            0.5 + card_width + spacing,
            top_y,
            card_width,
            card_height,
            value_color=self.theme.accent_color,
        )
        
        # EBITDA
        self._add_metric_card(
            slide,
            "EBITDA",
            format_currency(metrics.get('ebitda', 0)),
            0.5 + 2 * (card_width + spacing),
            top_y,
            card_width,
            card_height,
            value_color=self.theme.heading_level_two_color,
        )
        
        # Bottom row metrics
        bottom_y = top_y + card_height + 0.4
        
        # Net Income
        self._add_metric_card(
            slide,
            "Net Income",
            format_currency(metrics.get('net_income', 0)),
            0.5,
            bottom_y,
            card_width,
            card_height,
            value_color=self.theme.positive_color,
        )
        
        # EBITDA Margin
        self._add_metric_card(
            slide,
            "EBITDA Margin",
            format_percent(metrics.get('ebitda_margin', 0)),
            0.5 + card_width + spacing,
            bottom_y,
            card_width,
            card_height,
            value_color=self.theme.heading_level_two_color,
        )
        
        # Growth Rate - Dynamic color based on growth
        growth = metrics.get('revenue_growth', 0)
        growth_color = self.theme.positive_color if growth > 0 else self.theme.negative_color
        self._add_metric_card(
            slide, "Revenue Growth", format_percent(growth),
            0.5 + 2 * (card_width + spacing), bottom_y, card_width, card_height, value_color=growth_color
        )
    
    def _add_revenue_breakdown_slide(self, title: str, revenue_data: Dict[str, Any]):
        """Create a revenue breakdown slide"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
            self.prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.theme.heading_level_one_color
        title_bar.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.name = self.theme.font_family_heading
        title_para.font.color.rgb = self.theme.surface_color
        title_para.alignment = PP_ALIGN.LEFT
        
        # Revenue breakdown content
        top = Inches(1.8)
        
        def format_currency(value):
            if isinstance(value, (int, float)):
                if value >= 1_000_000:
                    return f"${value/1_000_000:.2f}M"
                elif value >= 1_000:
                    return f"${value/1_000:.2f}K"
                return f"${value:,.2f}"
            return str(value)
        
        # Total Revenue highlight
        total_revenue = revenue_data.get('total', revenue_data.get('revenue', 0))
        total_box = slide.shapes.add_textbox(
            Inches(0.5), top, Inches(9), Inches(1)
        )
        total_frame = total_box.text_frame
        total_frame.text = f"Total Revenue: {format_currency(total_revenue)}"
        total_para = total_frame.paragraphs[0]
        total_para.font.size = Pt(32)
        total_para.font.bold = True
        total_para.font.name = self.theme.font_family_heading
        total_para.font.color.rgb = self.theme.heading_level_one_color
        total_para.alignment = PP_ALIGN.CENTER
        
        # Breakdown items
        breakdown_items = revenue_data.get('breakdown', {})
        current_top = top + Inches(1.3)
        row_height = Inches(0.5)
        
        for i, (category, amount) in enumerate(breakdown_items.items()):
            # Row background
            if i % 2 == 0:
                row_bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(1), current_top,
                    Inches(8), row_height
                )
                row_bg.fill.solid()
                row_bg.fill.fore_color.rgb = self.theme.surface_color
                row_bg.line.fill.background()
            
            # Category name
            cat_box = slide.shapes.add_textbox(
                Inches(1.2), current_top + Inches(0.1), Inches(5), row_height - Inches(0.2)
            )
            cat_frame = cat_box.text_frame
            cat_frame.text = category
            cat_para = cat_frame.paragraphs[0]
            cat_para.font.size = Pt(16)
            cat_para.font.bold = True
            cat_para.font.name = self.theme.font_family_body
            cat_para.font.color.rgb = self.theme.text_color
            
            # Amount
            amount_box = slide.shapes.add_textbox(
                Inches(6), current_top + Inches(0.1), Inches(2.8), row_height - Inches(0.2)
            )
            amount_frame = amount_box.text_frame
            amount_frame.text = format_currency(amount)
            amount_para = amount_frame.paragraphs[0]
            amount_para.font.size = Pt(16)
            amount_para.font.bold = True
            amount_para.font.name = self.theme.font_family_heading
            amount_para.font.color.rgb = self.theme.accent_color
            amount_para.alignment = PP_ALIGN.RIGHT
            
            current_top += row_height
    
    def generate_quarterly_report(
        self,
        company_data: Dict[str, Any],
        output_path: str,
        custom_slides: Optional[List[Dict[str, Any]]] = None,
    ):
        """
        Generate a complete quarterly report presentation
        
        Args:
            company_data: Dictionary containing:
                - company_name: Company name
                - quarter: Quarter identifier (e.g., "Q4 2024")
                - revenue: Revenue amount
                - cogs: Cost of Goods Sold
                - gross_profit: Gross Profit
                - gross_margin: Gross Margin percentage
                - opex: Operating Expenses
                - ebitda: EBITDA
                - ebitda_margin: EBITDA Margin percentage
                - depreciation: Depreciation & Amortization
                - ebit: EBIT
                - interest: Interest Expense
                - net_income: Net Income
                - revenue_growth: Revenue Growth percentage
                - revenue_breakdown: Dict of revenue by category (optional)
            output_path: Path to save the presentation
            custom_slides: Optional list of dictionaries with custom slide specs
        """
        # Intro slide
        company_name = company_data.get('company_name', 'Company')
        quarter = company_data.get('quarter', 'Q1 2024')
        self._add_title_slide(
            f"{company_name}",
            f"{quarter} Quarterly Report"
        )
        
        # Key Metrics Overview
        metrics = {
            'revenue': company_data.get('revenue', 0),
            'gross_margin': company_data.get('gross_margin', 0),
            'ebitda': company_data.get('ebitda', 0),
            'ebitda_margin': company_data.get('ebitda_margin', 0),
            'net_income': company_data.get('net_income', 0),
            'revenue_growth': company_data.get('revenue_growth', 0),
        }
        self._add_key_metrics_slide("Key Financial Metrics", metrics)
        
        # P&L Statement
        pnl_data = {
            'revenue': company_data.get('revenue', 0),
            'cogs': company_data.get('cogs', 0),
            'gross_profit': company_data.get('gross_profit', 0),
            'opex': company_data.get('opex', 0),
            'ebitda': company_data.get('ebitda', 0),
            'depreciation': company_data.get('depreciation', 0),
            'ebit': company_data.get('ebit', 0),
            'interest': company_data.get('interest', 0),
            'net_income': company_data.get('net_income', 0),
        }
        self._add_pnl_slide("Profit & Loss Statement", pnl_data)
        
        # Revenue Breakdown (if provided)
        if company_data.get('revenue_breakdown'):
            revenue_data = {
                'total': company_data.get('revenue', 0),
                'breakdown': company_data.get('revenue_breakdown', {})
            }
            self._add_revenue_breakdown_slide("Revenue Breakdown", revenue_data)
        
        # Custom slides appended at the end
        if custom_slides:
            for slide_spec in custom_slides:
                title = slide_spec.get("title", "Additional Insights")
                content = slide_spec.get("content", "")
                self._add_content_slide(title, content)
        
        # Save presentation
        self.save(output_path)
        return output_path
    
    def save(self, filename: str):
        """Save the presentation to a file"""
        self.prs.save(filename)


def generate_pptx_from_json(
    json_data: str,
    output_path: str,
    theme: Optional[PresentationTheme] = None,
) -> str:
    """
    Convenience function to generate PPTX from JSON string
    
    Args:
        json_data: JSON string containing slides data
        output_path: Path where the PPTX file should be saved
    
    Returns:
        Path to the saved file
    """
    slides_data = json.loads(json_data)
    generator = SlideGenerator(theme=theme)
    generator.generate_from_json(slides_data)
    generator.save(output_path)
    return output_path


def generate_pptx_from_dict(
    slides_data: List[Dict[str, Any]],
    output_path: str,
    theme: Optional[PresentationTheme] = None,
) -> str:
    """
    Convenience function to generate PPTX from Python dictionary
    
    Args:
        slides_data: List of dictionaries containing slide data
        output_path: Path where the PPTX file should be saved
    
    Returns:
        Path to the saved file
    """
    generator = SlideGenerator(theme=theme)
    generator.generate_from_json(slides_data)
    generator.save(output_path)
    return output_path


def generate_quarterly_report_template(
    company_data: Dict[str, Any],
    output_path: str,
    theme: Optional[PresentationTheme] = None,
    custom_slides: Optional[List[Dict[str, Any]]] = None,
) -> str:
    """
    Convenience function to generate quarterly report from company data
    
    Args:
        company_data: Dictionary containing financial data (see generate_quarterly_report docstring)
        output_path: Path where the PPTX file should be saved
        theme: Optional custom PresentationTheme instance
    
    Returns:
        Path to the saved file
    
    Example:
        company_data = {
            "company_name": "TechCorp Inc.",
            "quarter": "Q4 2024",
            "revenue": 10_000_000,
            "cogs": 4_000_000,
            "gross_profit": 6_000_000,
            "gross_margin": 60.0,
            "opex": 3_000_000,
            "ebitda": 3_000_000,
            "ebitda_margin": 30.0,
            "depreciation": 500_000,
            "ebit": 2_500_000,
            "interest": 200_000,
            "net_income": 2_300_000,
            "revenue_growth": 25.5,
            "revenue_breakdown": {
                "Product Sales": 6_000_000,
                "Services": 3_000_000,
                "Licensing": 1_000_000
            }
        }
        generate_quarterly_report_template(company_data, "q4_report.pptx")
    """
    generator = SlideGenerator(theme=theme)
    generator.generate_quarterly_report(company_data, output_path, custom_slides=custom_slides)
    return output_path


# Example usage
if __name__ == "__main__":
    # Example JSON structure
    example_slides = [
        {
            "type": "title",
            "title": "Welcome to Finstory",
            "subtitle": "Financial Storytelling Made Easy"
        },
        {
            "type": "content",
            "title": "Overview",
            "content": [
                "Financial data analysis and visualization",
                "Automated report generation",
                "Beautiful presentation creation",
                "Data-driven insights"
            ]
        },
        {
            "type": "content",
            "title": "Key Features",
            "content": {
                "type": "two_column",
                "left": [
                    "Real-time data processing",
                    "Advanced analytics",
                    "Custom templates"
                ],
                "right": [
                    "Export to multiple formats",
                    "Collaborative editing",
                    "Secure data handling"
                ]
            }
        },
        {
            "type": "content",
            "title": "Statistics",
            "content": {
                "Total Users": "10,000+",
                "Data Points": "1M+",
                "Reports Generated": "50K+",
                "Accuracy": "99.9%"
            }
        }
    ]
    
    # Generate presentation
    generator = SlideGenerator()
    generator.generate_from_json(example_slides)
    generator.save("example_presentation.pptx")
    print("Presentation generated: example_presentation.pptx")
    
    # Example quarterly report
    example_company_data = {
        "company_name": "TechCorp Inc.",
        "quarter": "Q4 2024",
        "revenue": 10_000_000,
        "cogs": 4_000_000,
        "gross_profit": 6_000_000,
        "gross_margin": 60.0,
        "opex": 3_000_000,
        "ebitda": 3_000_000,
        "ebitda_margin": 30.0,
        "depreciation": 500_000,
        "ebit": 2_500_000,
        "interest": 200_000,
        "net_income": 2_300_000,
        "revenue_growth": 25.5,
        "revenue_breakdown": {
            "Product Sales": 6_000_000,
            "Services": 3_000_000,
            "Licensing": 1_000_000
        }
    }
    
    # Generate quarterly report
    generate_quarterly_report_template(example_company_data, "quarterly_report_example.pptx")
    print("Quarterly report generated: quarterly_report_example.pptx")