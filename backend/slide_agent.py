"""
LangGraph-powered agent that reads a financial prompt, decides on a color scheme,
and generates a quarterly slide deck using SlideGenerator.
"""

from __future__ import annotations

import json
import logging
import os
import re
from dataclasses import dataclass, replace
from pathlib import Path
from typing import Any, Dict, List, Optional, TypedDict

from langgraph.graph import END, StateGraph
from pptx.dml.color import RGBColor
from dotenv import load_dotenv

load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
try:
    import google.generativeai as genai
except ImportError:  # pragma: no cover - optional dependency
    genai = None  # type: ignore

try:  # Works whether executed as module or script
    from .slidegen import SlideGenerator, PresentationTheme, DEFAULT_THEME
except ImportError:  # pragma: no cover
    from slidegen import SlideGenerator, PresentationTheme, DEFAULT_THEME  # type: ignore


# ---------------------------------------------------------------------------
# Helpers for extracting numerical signals from the user prompt
# ---------------------------------------------------------------------------

CURRENCY_HINTS = {
    "billion": 1_000_000_000,
    "bn": 1_000_000_000,
    "b": 1_000_000_000,
    "million": 1_000_000,
    "mm": 1_000_000,
    "m": 1_000_000,
    "thousand": 1_000,
    "k": 1_000,
}


METRIC_KEYWORDS = {
    "revenue": ["revenue", "sales", "top line"],
    "cogs": ["cogs", "cost of goods", "cost of revenue"],
    "gross_profit": ["gross profit"],
    "gross_margin": ["gross margin"],
    "opex": ["opex", "operating expense", "operating cost"],
    "ebitda": ["ebitda"],
    "ebitda_margin": ["ebitda margin"],
    "depreciation": ["depreciation", "amortization"],
    "ebit": ["ebit"],
    "interest": ["interest expense"],
    "net_income": ["net income", "bottom line", "earnings"],
    "revenue_growth": ["revenue growth", "growth rate"],
}


def parse_currency_token(token: str) -> Optional[float]:
    if not token:
        return None
    cleaned = token.lower().replace(",", "").replace("$", "").replace("€", "").replace("£", "").strip()
    multiplier = 1.0
    for hint, factor in CURRENCY_HINTS.items():
        if hint in cleaned:
            cleaned = cleaned.replace(hint, "")
            multiplier = factor
            break
    match = re.search(r"[-+]?\d*\.?\d+", cleaned)
    if not match:
        return None
    value = float(match.group())
    return value * multiplier


def parse_percentage_token(token: str) -> Optional[float]:
    if not token:
        return None
    match = re.search(r"[-+]?\d*\.?\d+", token.replace("%", ""))
    if not match:
        return None
    return float(match.group())


def extract_metric(prompt: str, keywords: list[str], percent: bool = False) -> Optional[float]:
    for keyword in keywords:
        pattern = rf"{re.escape(keyword)}[^\d$€£%]*([\$€£]?\s?[-+]?\d[\d,\.]*(?:\s?(?:billion|million|thousand)|[MBKmbk%]))"
        match = re.search(pattern, prompt, re.IGNORECASE)
        if match:
            token = match.group(1)
            return parse_percentage_token(token) if percent else parse_currency_token(token)
    return None


def extract_company_name(prompt: str) -> str:
    explicit = re.search(r"company\s*[:\-]\s*([A-Za-z0-9 &]+)", prompt, re.IGNORECASE)
    if explicit:
        return explicit.group(1).strip().title()
    match = re.search(r"for\s+([A-Z][A-Za-z0-9 &]+)", prompt)
    if match:
        return match.group(1).strip().title()
    return "Finstory Client"


def extract_quarter(prompt: str) -> str:
    match = re.search(r"(Q[1-4])\s*['’]?\s*(\d{2,4})", prompt, re.IGNORECASE)
    if match:
        quarter = match.group(1).upper()
        year_fragment = match.group(2)
        year = int(year_fragment)
        if year < 100:
            year += 2000
        return f"{quarter} {year}"
    return "Q1 2024"


def extract_revenue_breakdown(prompt: str) -> Dict[str, float]:
    breakdown: Dict[str, float] = {}
    pattern = re.compile(r"([A-Za-z &]{3,40}):\s*([\$€£]?\s?[-+]?\d[\d,\.]*(?:\s?(?:billion|million|thousand)|[MBKmbk]))")
    for label, token in pattern.findall(prompt):
        key = label.strip().title()
        canonical = key.lower()
        if any(canonical.startswith(metric) for metric in METRIC_KEYWORDS):
            continue
        value = parse_currency_token(token)
        if value is not None:
            breakdown[key] = value
    return breakdown


def infer_company_data(prompt: str) -> Dict[str, Any]:
    prompt_lower = prompt.lower()
    data: Dict[str, Any] = {
        "company_name": extract_company_name(prompt),
        "quarter": extract_quarter(prompt),
    }
    for metric, keywords in METRIC_KEYWORDS.items():
        percent = metric.endswith("margin") or metric == "revenue_growth"
        data[metric] = extract_metric(prompt, keywords, percent=percent)

    # Derive missing values where possible
    revenue = data.get("revenue")
    cogs = data.get("cogs")
    if data.get("gross_profit") is None and revenue is not None and cogs is not None:
        data["gross_profit"] = revenue - cogs

    if data.get("gross_margin") is None and revenue:
        gross_profit = data.get("gross_profit")
        if gross_profit:
            data["gross_margin"] = (gross_profit / revenue) * 100

    if data.get("ebitda") is None:
        gross_profit = data.get("gross_profit")
        opex = data.get("opex")
        if gross_profit is not None and opex is not None:
            data["ebitda"] = gross_profit - opex

    if data.get("ebitda_margin") is None and data.get("ebitda") and revenue:
        data["ebitda_margin"] = (data["ebitda"] / revenue) * 100

    if data.get("ebit") is None and data.get("ebitda") is not None:
        depreciation = data.get("depreciation") or 0
        data["ebit"] = data["ebitda"] - depreciation

    if data.get("net_income") is None and data.get("ebit") is not None:
        interest = data.get("interest") or 0
        data["net_income"] = data["ebit"] - interest

    # Revenue breakdown categories
    data["revenue_breakdown"] = extract_revenue_breakdown(prompt)

    # Defaults to avoid None in downstream logic
    numeric_defaults = [
        "revenue",
        "cogs",
        "gross_profit",
        "gross_margin",
        "opex",
        "ebitda",
        "ebitda_margin",
        "depreciation",
        "ebit",
        "interest",
        "net_income",
        "revenue_growth",
    ]
    for field in numeric_defaults:
        if data.get(field) is None:
            data[field] = 0.0

    return data


# ---------------------------------------------------------------------------
# Theme selection helpers
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class ThemePreset:
    name: str
    keywords: tuple[str, ...]
    theme: PresentationTheme


THEME_LIBRARY: tuple[ThemePreset, ...] = (
    ThemePreset(
        name="polar",
        keywords=("luxury", "premium", "executive", "minimal"),
        theme=PresentationTheme(
            background_color=RGBColor(241, 245, 249),
            heading_level_one_color=RGBColor(15, 76, 129),
            heading_level_two_color=RGBColor(59, 130, 246),
            text_color=RGBColor(30, 41, 59),
            surface_color=RGBColor(255, 255, 255),
            accent_color=RGBColor(79, 70, 229),
            positive_color=RGBColor(34, 197, 94),
            negative_color=RGBColor(239, 68, 68),
        ),
    ),
    ThemePreset(
        name="fjord",
        keywords=("green", "sustainable", "climate", "eco", "clean"),
        theme=PresentationTheme(
            background_color=RGBColor(236, 253, 245),
            heading_level_one_color=RGBColor(5, 122, 85),
            heading_level_two_color=RGBColor(14, 165, 233),
            text_color=RGBColor(15, 63, 42),
            surface_color=RGBColor(255, 255, 255),
            accent_color=RGBColor(59, 130, 246),
            positive_color=RGBColor(22, 163, 74),
            negative_color=RGBColor(225, 29, 72),
        ),
    ),
    ThemePreset(
        name="azure",
        keywords=("tech", "digital", "ai", "innovation", "future"),
        theme=PresentationTheme(
            background_color=RGBColor(239, 246, 255),
            heading_level_one_color=RGBColor(23, 37, 84),
            heading_level_two_color=RGBColor(56, 189, 248),
            text_color=RGBColor(30, 41, 59),
            surface_color=RGBColor(255, 255, 255),
            accent_color=RGBColor(59, 130, 246),
            positive_color=RGBColor(34, 197, 94),
            negative_color=RGBColor(248, 113, 113),
        ),
    ),
)


def choose_theme(prompt: str) -> PresentationTheme:
    prompt_lower = prompt.lower()
    for preset in THEME_LIBRARY:
        if any(keyword in prompt_lower for keyword in preset.keywords):
            return preset.theme
    return DEFAULT_THEME


STYLE_OVERRIDES = {
    "minimal": {
        "font_family_heading": "Helvetica Neue",
        "font_family_body": "Helvetica",
        "accent_color": RGBColor(71, 85, 105),
    },
    "playful": {
        "font_family_heading": "Poppins",
        "font_family_body": "Nunito",
        "accent_color": RGBColor(236, 72, 153),
        "heading_level_two_color": RGBColor(244, 114, 182),
    },
    "bold": {
        "font_family_heading": "Montserrat",
        "font_family_body": "Inter",
        "accent_color": RGBColor(59, 130, 246),
    },
    "corporate": {
        "font_family_heading": "IBM Plex Sans",
        "font_family_body": "IBM Plex Sans",
        "accent_color": RGBColor(15, 118, 110),
    },
}


def apply_style_keywords(base_theme: PresentationTheme, prompt: str) -> PresentationTheme:
    theme = base_theme
    prompt_lower = prompt.lower()
    overrides: Dict[str, Any] = {}
    for keyword, update in STYLE_OVERRIDES.items():
        if keyword in prompt_lower:
            overrides.update(update)
    if overrides:
        theme = replace(theme, **overrides)
    return theme


def determine_theme(prompt: str) -> PresentationTheme:
    base = choose_theme(prompt)
    return apply_style_keywords(base, prompt)


logger = logging.getLogger(__name__)

SLIDE_TOOLBOX = [
    {
        "name": "hero",
        "description": "Headline slide with big title + subtitle describing the moment.",
        "args": {"title": "string", "subtitle": "string"},
    },
    {
        "name": "vision",
        "description": "Tone/vision slide with inspirational bullets.",
        "args": {"title": "string", "bullets": ["string", "..."]},
    },
    {
        "name": "kpi_grid",
        "description": "Key financials (Revenue, Gross Margin, EBITDA, Net Income, Growth).",
        "args": {"title": "string", "notes": "optional string"},
    },
    {
        "name": "story",
        "description": "Narrative paragraph to explain context.",
        "args": {"title": "string", "text": "string"},
    },
    {
        "name": "two_column",
        "description": "Split bullets into left/right themes.",
        "args": {"title": "string", "left": ["string"], "right": ["string"]},
    },
    {
        "name": "bullets",
        "description": "Standard bullet slide.",
        "args": {"title": "string", "bullets": ["string"]},
    },
]

def tool_instruction() -> str:
    lines = ["Available slide tools:"]
    for tool in SLIDE_TOOLBOX:
        lines.append(f"- {tool['name']}: {tool['description']} Args: {tool['args']}")
    return "\n".join(lines)


class GeminiPlanner:
    """Wrapper around Gemini models to produce narrative + slide suggestions."""

    def __init__(
        self,
        model_name: str = "gemini-2.5-flash",
        api_key: Optional[str] = GEMINI_API_KEY,
        temperature: float = 0.25,
    ):
        if genai is None:
            raise RuntimeError(
                "google-generativeai is not installed. Install backend requirements."
            )
        key = api_key or os.environ.get("GEMINI_API_KEY")
        if not key:
            raise RuntimeError("Set GEMINI_API_KEY to enable Gemini-powered planning.")
        genai.configure(api_key=key)
        self.model = genai.GenerativeModel(model_name)
        self.generation_config = {
            "temperature": temperature,
            "top_p": 0.8,
            "top_k": 40,
            "response_mime_type": "application/json",
        }

    def plan(self, prompt: str, company_data: Dict[str, Any]) -> Dict[str, Any]:
        payload = {
            "prompt": prompt,
            "company_data": company_data,
        }
        instructions = (
            "You are a financial storytelling assistant. "
            "Analyze the prompt and company_data JSON below. "
            f"{tool_instruction()}\n"
            "Keep slide titles simple (2-4 words) and under 40 characters.\n"
            "Return ONLY JSON with the schema:\n"
            "{\n"
            '  "narrative_highlights": ["sentence", ...],\n'
            '  "slides": [\n'
            '      {"tool": "<one of the tools>", "title": "string", ... args ... }\n'
            "  ],\n"
            '  "theme_hint": "optional descriptive words"\n'
            "}\n"
            "Bullets should incorporate financial insights from company_data where possible."
        )
        response = self.model.generate_content(
            [instructions, json.dumps(payload)],
            generation_config=self.generation_config,
        )
        text = (response.text or "").strip()
        if not text:
            raise ValueError("Gemini returned empty response")
        try:
            return json.loads(text)
        except json.JSONDecodeError as exc:
            raise ValueError(f"Failed to decode Gemini JSON: {exc}") from exc


CUSTOM_TOPIC_KEYWORDS = {
    "Product Strategy": ("product", "innovation", "roadmap", "feature"),
    "Customer Highlights": ("customer", "client", "user", "partner"),
    "Market Opportunity": ("market", "industry", "opportunity", "demand"),
    "Risks & Mitigation": ("risk", "challenge", "concern", "threat"),
    "Capital Plan": ("capital", "funding", "cash", "runway"),
    "Sustainability & Impact": ("esg", "sustainability", "impact", "green", "climate"),
}


TITLE_FALLBACK = "Update"
MAX_TITLE_LENGTH = 44


def normalize_title_text(raw_title: Optional[str], fallback: str = TITLE_FALLBACK) -> str:
    title = (raw_title or "").strip()
    if not title:
        title = fallback
    # Keep only the first sentence/phrase to avoid long descriptions
    primary = re.split(r"[.!?\n]", title, 1)[0].strip()
    if primary:
        title = primary
    if len(title) > MAX_TITLE_LENGTH:
        truncated = title[:MAX_TITLE_LENGTH].rstrip()
        last_space = truncated.rfind(" ")
        if last_space > 20:
            truncated = truncated[:last_space]
        title = truncated.rstrip(" -–—,:;")
    return title or fallback


def split_sentences(prompt: str) -> List[str]:
    normalized = prompt.replace("\r", " ")
    parts = re.split(r'(?<=[.!?])\s+|\n', normalized)
    sentences: List[str] = []
    for part in parts:
        cleaned = part.strip(" -•\t")
        if len(cleaned) >= 8:
            sentences.append(cleaned)
    return sentences


def parse_explicit_slide_requests(prompt: str) -> List[Dict[str, Any]]:
    slides: List[Dict[str, Any]] = []
    for line in prompt.splitlines():
        if not line.lower().strip().startswith("slide"):
            continue
        _, _, rest = line.partition(":")
        rest = rest.strip()
        if not rest:
            continue
        title, bullets_str = (rest.split("-", 1) + [""])[:2]
        bullets = [item.strip() for item in re.split(r"[;|]", bullets_str) if item.strip()]
        if bullets:
            slides.append({
                "type": "content",
                "title": normalize_title_text(title.strip().title()),
                "content": {"type": "bullets", "items": bullets[:6]},
            })
    return slides


def extract_topic_slides(prompt: str) -> List[Dict[str, Any]]:
    sentences = split_sentences(prompt)
    slides: List[Dict[str, Any]] = []
    used_sentences = set()
    for title, keywords in CUSTOM_TOPIC_KEYWORDS.items():
        matched: List[str] = []
        for sentence in sentences:
            if sentence in used_sentences:
                continue
            if any(keyword in sentence.lower() for keyword in keywords):
                matched.append(sentence)
                used_sentences.add(sentence)
            if len(matched) == 4:
                break
        if matched:
            slides.append({
                "type": "content",
                "title": normalize_title_text(title),
                "content": {"type": "bullets", "items": matched},
            })
    return slides


def format_currency_short(value: Optional[float]) -> str:
    if value is None:
        return "N/A"
    abs_value = abs(value)
    if abs_value >= 1_000_000_000:
        return f"${value/1_000_000_000:.1f}B"
    if abs_value >= 1_000_000:
        return f"${value/1_000_000:.1f}M"
    if abs_value >= 1_000:
        return f"${value/1_000:.1f}K"
    return f"${value:,.0f}"


def format_percent(value: Optional[float]) -> str:
    if value is None:
        return "N/A"
    return f"{value:.1f}%"


def build_financial_snapshot_slide(company_data: Dict[str, Any]) -> Dict[str, Any]:
    bullets = [
        f"Revenue: {format_currency_short(company_data.get('revenue'))}",
        f"Gross Margin: {format_percent(company_data.get('gross_margin'))}",
        f"EBITDA: {format_currency_short(company_data.get('ebitda'))} "
        f"({format_percent(company_data.get('ebitda_margin'))})",
        f"Net Income: {format_currency_short(company_data.get('net_income'))}",
    ]
    return {
        "type": "content",
        "title": normalize_title_text("Financial Snapshot"),
        "content": {"type": "bullets", "items": bullets},
    }


def build_prompt_summary_slide(prompt: str) -> Dict[str, Any]:
    sentences = split_sentences(prompt)[:4]
    if not sentences:
        sentences = ["Custom insights requested via prompt."]
    return {
        "type": "content",
        "title": normalize_title_text("Narrative Highlights"),
        "content": {"type": "bullets", "items": sentences},
    }


def build_custom_slides(prompt: str, company_data: Dict[str, Any]) -> List[Dict[str, Any]]:
    slides: List[Dict[str, Any]] = [build_financial_snapshot_slide(company_data)]
    slides.extend(parse_explicit_slide_requests(prompt))
    topic_slides = extract_topic_slides(prompt)
    for slide in topic_slides:
        if slide["title"] not in (existing["title"] for existing in slides):
            slides.append(slide)
    if len(slides) == 1:
        slides.append(build_prompt_summary_slide(prompt))
    return slides[:6]


# ---------------------------------------------------------------------------
# LangGraph state + workflow
# ---------------------------------------------------------------------------


class AgentState(TypedDict, total=False):
    prompt: str
    requested_output_path: str
    company_data: Dict[str, Any]
    slide_plan: List[Dict[str, Any]]
    theme: PresentationTheme
    llm_theme_hint: str
    llm_narrative: List[str]
    report_path: str


class SlideAgent:
    """
    LangGraph agent that:
    1. Analyzes the prompt for financial data.
    2. Optionally calls Gemini to craft narratives + slide ideas.
    3. Chooses a theme based on intent/keywords (and LLM hints).
    4. Builds a quarterly report slide deck.
    """

    def __init__(
        self,
        default_output_path: str = "agent_quarterly_report.pptx",
        use_gemini: bool = True,
        gemini_model: str = "gemini-2.5-flash",
        gemini_api_key: Optional[str] = None,
    ):
        self.default_output_path = default_output_path
        self.planner: Optional[GeminiPlanner] = None
        if use_gemini:
            try:
                self.planner = GeminiPlanner(gemini_model, gemini_api_key)
                logger.info("Gemini planner enabled (model=%s)", gemini_model)
            except Exception as exc:  # pragma: no cover - optional runtime
                logger.warning("Gemini planner disabled: %s", exc)
        self.graph = self._build_graph()

    # Graph construction -----------------------------------------------------
    def _build_graph(self):
        workflow = StateGraph(AgentState)
        workflow.add_node("analyze_prompt", self._analyze_prompt_node)
        workflow.add_node("llm_generate", self._llm_generate_node)
        workflow.add_node("select_theme", self._select_theme_node)
        workflow.add_node("build_deck", self._build_deck_node)
        workflow.set_entry_point("analyze_prompt")
        workflow.add_edge("analyze_prompt", "llm_generate")
        workflow.add_edge("llm_generate", "select_theme")
        workflow.add_edge("select_theme", "build_deck")
        workflow.add_edge("build_deck", END)
        return workflow.compile()

    # Graph nodes ------------------------------------------------------------
    def _analyze_prompt_node(self, state: AgentState) -> AgentState:
        prompt = state.get("prompt", "")
        company_data = infer_company_data(prompt)
        fallback_plan = build_custom_slides(prompt, company_data)
        return {"company_data": company_data, "slide_plan": fallback_plan}

    def _llm_generate_node(self, state: AgentState) -> AgentState:
        if not self.planner:
            return {}
        prompt = state.get("prompt", "")
        company_data = state.get("company_data", {})
        try:
            plan = self.planner.plan(prompt, company_data)
        except Exception as exc:  # pragma: no cover - relies on external API
            logger.warning("Gemini plan failed: %s", exc)
            return {}

        updates: AgentState = {}
        llm_slides = self._normalize_llm_slides(plan.get("slides", []), company_data)

        if llm_slides:
            heuristic_slides = state.get("slide_plan", [])
            merged = self._merge_slides(llm_slides, heuristic_slides)
            updates["slide_plan"] = merged

        narrative = plan.get("narrative_highlights")
        if isinstance(narrative, list) and narrative:
            updates["llm_narrative"] = narrative
            narrative_slide = {
                "title": normalize_title_text("Narrative Highlights"),
                "content": {"type": "bullets", "items": narrative[:6]},
            }
            merged = self._merge_slides(
                [narrative_slide],
                updates.get("slide_plan") or state.get("slide_plan", []),
            )
            updates["slide_plan"] = merged

        theme_hint = plan.get("theme_hint")
        if isinstance(theme_hint, str) and theme_hint.strip():
            updates["llm_theme_hint"] = theme_hint.strip()

        return updates

    def _select_theme_node(self, state: AgentState) -> AgentState:
        prompt = state.get("prompt", "")
        theme_hint = state.get("llm_theme_hint")
        if theme_hint:
            prompt = f"{prompt}\nPreferred vibe: {theme_hint}"
        theme = determine_theme(prompt)
        return {"theme": theme}

    def _build_deck_node(self, state: AgentState) -> AgentState:
        company_data = state["company_data"]
        theme = state["theme"]
        slide_plan = state.get("slide_plan") or build_custom_slides("", company_data)
        output_path = state.get("requested_output_path") or self.default_output_path
        output_path = str(Path(output_path).expanduser().resolve())
        generator = SlideGenerator(theme=theme)
        generator.generate_from_json(slide_plan)
        generator.save(output_path)
        return {"report_path": output_path}

    # Helpers ----------------------------------------------------------------
    @staticmethod
    def _normalize_llm_slides(slides: Any, company_data: Dict[str, Any]) -> List[Dict[str, Any]]:
        normalized: List[Dict[str, Any]] = []
        if not isinstance(slides, list):
            return normalized
        for raw in slides:
            if not isinstance(raw, dict):
                continue
            slide = SlideAgent._tool_to_slide(raw, company_data)
            if slide:
                normalized.append(slide)
        return normalized[:6]

    @staticmethod
    def _tool_to_slide(call: Dict[str, Any], company_data: Dict[str, Any]) -> Optional[Dict[str, Any]]:
        tool = (call.get("tool") or call.get("type") or "").lower()
        fallback_title = company_data.get("company_name", TITLE_FALLBACK)
        title = normalize_title_text(call.get("title"), fallback=fallback_title)

        if tool == "hero":
            subtitle = call.get("subtitle") or company_data.get("quarter", "")
            return {"type": "title", "title": title, "subtitle": subtitle}

        if tool == "vision":
            bullets = call.get("bullets") or call.get("items") or []
            if not bullets:
                return None
            return {"type": "content", "title": title, "content": {"type": "bullets", "items": bullets[:8]}}

        if tool == "story":
            text = call.get("text") or call.get("body")
            if not text:
                return None
            return {"type": "content", "title": title, "content": {"type": "text", "text": text}}

        if tool == "two_column":
            left = call.get("left") or []
            right = call.get("right") or []
            if not left and not right:
                return None
            return {
                "type": "content",
                "title": title,
                "content": {"type": "two_column", "left": left[:6], "right": right[:6]},
            }

        if tool == "kpi_grid":
            bullets = [
                f"Revenue: {format_currency_short(company_data.get('revenue'))}",
                f"Gross Margin: {format_percent(company_data.get('gross_margin'))}",
                f"EBITDA: {format_currency_short(company_data.get('ebitda'))} "
                f"({format_percent(company_data.get('ebitda_margin'))})",
                f"Net Income: {format_currency_short(company_data.get('net_income'))}",
                f"Revenue Growth: {format_percent(company_data.get('revenue_growth'))}",
            ]
            if call.get("notes"):
                bullets.append(call["notes"])
            return {"type": "content", "title": title or "Key Metrics", "content": {"type": "bullets", "items": bullets}}

        if tool == "bullets":
            bullets = call.get("bullets") or call.get("items") or []
            if not bullets:
                return None
            return {"type": "content", "title": title, "content": {"type": "bullets", "items": bullets[:8]}}

        # fallback
        bullets = call.get("bullets") or call.get("items") or []
        if bullets:
            return {"type": "content", "title": title, "content": {"type": "bullets", "items": bullets[:6]}}
        text = call.get("text")
        if text:
            return {"type": "content", "title": title, "content": {"type": "text", "text": text}}
        return None

    @staticmethod
    def _merge_slides(priority: List[Dict[str, Any]], fallback: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        merged: List[Dict[str, Any]] = []
        seen_titles: set[str] = set()

        def add_slide(slide: Dict[str, Any]):
            title = normalize_title_text(slide.get("title"))
            slide = {**slide, "title": title}
            key = title.lower()
            if key in seen_titles:
                return
            seen_titles.add(key)
            merged.append(slide)

        for slide in priority:
            add_slide(slide)
        for slide in fallback:
            add_slide(slide)
        return merged[:8]

    # Public API -------------------------------------------------------------
    def run(self, prompt: str, output_path: Optional[str] = None) -> AgentState:
        initial_state: AgentState = {
            "prompt": prompt,
            "requested_output_path": output_path or self.default_output_path,
        }
        return self.graph.invoke(initial_state)


# ---------------------------------------------------------------------------
# Usage example
# ---------------------------------------------------------------------------


if __name__ == "__main__":
    agent = SlideAgent()
    sample_prompt = """
    Create an exclusive, premium investor update for Aurora Luxe Holdings for Q4 2024.
    Revenue hit $24M with COGS of $9.5M, gross margin ~60%, opex $7M, EBITDA $7.5M,
    depreciation $1.1M, interest expense $0.4M, and net income around $6M.
    Revenue growth was 18%. Product Sales: $14M, Services: $6M, Partnerships: $4M.
    """
    result = agent.run(sample_prompt, output_path="aurora_q_agent.pptx")
    print(f"Generated report at {result['report_path']}")

